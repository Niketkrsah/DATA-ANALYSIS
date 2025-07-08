import os
import re
import json
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

plt.rcParams["font.family"] = "Segoe UI Emoji"
sns.set(style="whitegrid")

def run_anr_analysis(input_csv, output_dir):
    def ensure_dirs(base):
        img_dir = os.path.join(base, "anr_images")
        ppt_path = os.path.join(base, "anr_report.pptx")
        json_path = os.path.join(base, "anr_summary.json")

        if os.path.exists(img_dir):
            for f in os.listdir(img_dir):
                if f.endswith(".jpg") or f.endswith(".png"):
                    os.remove(os.path.join(img_dir, f))
        else:
            os.makedirs(img_dir)

        if os.path.exists(ppt_path):
            os.remove(ppt_path)

        if os.path.exists(json_path):
            os.remove(json_path)

        return img_dir, ppt_path, json_path

    def save_chart(fig, name, img_dir):
        path = os.path.join(img_dir, f"{name}.jpg")
        fig.savefig(path, dpi=200, bbox_inches="tight")
        plt.close(fig)
        return path

    def bar_plot(series, title, name, img_dir, horizontal=False):
        import textwrap
        series.index = ['\n'.join(textwrap.wrap(str(label), width=50)) for label in series.index]
        fig, ax = plt.subplots(figsize=(15, 6))
        if horizontal:
            series.plot(kind="barh", ax=ax, color="skyblue")
            ax.invert_yaxis()
            for p in ax.patches:
               val = p.get_width()
               ax.annotate(str(int(val)), (p.get_x() + val + 5, p.get_y() + p.get_height() / 2),
                            ha='left', va='center', fontsize=9)
        else:
            series.plot(kind="bar", ax=ax, color="steelblue")
            ax.set_xticklabels(ax.get_xticklabels(), rotation=45)
            for p in ax.patches:
                val = p.get_height()
                ax.annotate(str(int(val)), (p.get_x() + p.get_width() / 2, val + 1),
                            ha='center', fontsize=9)

        ax.set_title(title)
        # for p in ax.patches:
        #     val = p.get_width() if horizontal else p.get_height()
        #     ax.annotate(str(int(val)), (p.get_x() + p.get_width() / 2, val + 1),
        #                 ha='center', fontsize=9)

        fig.tight_layout(pad=3.0)
        return save_chart(fig, name, img_dir)

    def build_ppt(img_dir, ppt_path):
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for img in sorted(os.listdir(img_dir)):
            if img.endswith(".jpg"):
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(os.path.join(img_dir, img), 0, 0,
                                         width=prs.slide_width, height=prs.slide_height)
        prs.save(ppt_path)

    def extract_field(pattern, text):
        match = re.search(pattern, str(text))
        return match.group(1) if match else None

    def extract_json_field(value, field):
        try:
            data = json.loads(value.replace("'", '"'))
            return float(data.get(field, 0))
        except Exception:
            return None

    img_dir, ppt_path, json_path = ensure_dirs(output_dir)

    df = pd.read_csv(input_csv)
    df["ANR Activity"] = df["ANR Process Info"].apply(lambda x: extract_field(r'Activity:\s+([^\s]+)', x))
    df["Subject"] = df["ANR Process Info"].apply(lambda x: extract_field(r'Subject:\s+(.*?)\s+Build:', x))
    df["Event Hour"] = pd.to_datetime(df["Event Time"], errors='coerce').dt.hour
    df["Available Memory"] = df["Mi 4"].apply(lambda x: extract_json_field(x, "available"))

    bins = [0, 200, 400, 600, 800, 1000, 1200, 2000]
    labels = ['0–200 MB', '200–400 MB', '400–600 MB', '600–800 MB', '800–1000 MB', '1000–1200 MB', '1200–2000 MB']
    df["Memory Range"] = pd.cut(df["Available Memory"], bins=bins, labels=labels, right=False)

    summary = {}

    summary["Top ANR Activities"] = df["ANR Activity"].value_counts().nlargest(10).to_dict()
    bar_plot(pd.Series(summary["Top ANR Activities"]), "Top 10 ANR Activities", "top_anr_activities", img_dir, horizontal=True)

    summary["ANR by App Version"] = df["App version"].value_counts().nlargest(10).to_dict()
    bar_plot(pd.Series(summary["ANR by App Version"]), "ANR Count by App Version", "anr_by_version", img_dir)

    summary["ANR by State"] = df["Gis City"].value_counts().nlargest(10).to_dict()
    bar_plot(pd.Series(summary["ANR by State"]), "Top 10 States with ANRs", "anr_by_state", img_dir)

    hourly = df["Event Hour"].value_counts().sort_index()
    fig, ax = plt.subplots(figsize=(10, 6))
    hourly.plot(kind="line", marker='o', ax=ax, color="blue")
    ax.set_title("ANR Occurrences by Hour of Day")
    ax.set_xlabel("Hour")
    ax.set_ylabel("Count")
    ax.grid(True)
    fig.tight_layout(pad=3.0)
    save_chart(fig, "anr_by_hour", img_dir)
    summary["ANR by Hour"] = hourly.to_dict()

    mem_range_counts = df["Memory Range"].value_counts().sort_index()
    bar_plot(mem_range_counts, "ANR Count by Available Memory Range (Mi 4)", "mi4_available_memory_ranges", img_dir)
    summary["Available Memory Ranges"] = mem_range_counts.to_dict()

    df["Clean Subject"] = df["ANR Process Info"].apply(lambda text: extract_field(r'Subject:\s*(.*?)(?:\(|Build:|$)', text))
    vc = df["Clean Subject"].value_counts().nlargest(10)
    bar_plot(vc, "Top ANR ERROR TYPE (Simplified)", "anr_subject_simplified", img_dir, horizontal=True)
    summary["Simplified ANR Subjects"] = vc.to_dict()

    if "Di 9" in df.columns:
        top_stb_series = df["Di 9"].value_counts().nlargest(5)
        fig, ax = plt.subplots(figsize=(10, 6))
        top_stb_series.plot(kind="bar", ax=ax, color="royalblue")

        for i, (value, count) in enumerate(top_stb_series.items()):
            ax.text(i, count + 5, str(count), ha='center', fontsize=10)

        ax.set_title("Top 5 STB Series firmware")
        ax.set_ylabel("ANR Count")
        ax.set_xlabel("STB Series")
        ax.set_xticklabels(ax.get_xticklabels(), rotation=45)
        fig.tight_layout(pad=2.0)

        save_chart(fig, "top_5_stb_series_vertical", img_dir)
        summary["Top 5 STB Series"] = top_stb_series.to_dict()

    if "D Customer Product Type" in df.columns:
        df.rename(columns={"D Customer Product Type": "Customer Product Type"}, inplace=True)
        vc = df["Customer Product Type"].value_counts().nlargest(5)
        bar_plot(vc, "Top 5 Customer Product Types by ANR Count", "top_5_customer_product_types", img_dir)
        summary["Customer Product Types"] = vc.to_dict()

    focused_window_issues = df["Subject"].dropna().str.contains("does not have a focused window", case=False)
    count = int(focused_window_issues.sum())
    summary["No Focused Window"] = {
    "No Focused Window": count,
    "Other Subjects": int(len(df) - count)
    }
    labels = ['No Focused Window', 'Other Subjects']
    sizes = [count, len(df) - count]
    fig, ax = plt.subplots(figsize=(6, 6))
    ax.pie(sizes, labels=labels, autopct='%1.1f%%', colors=['tomato', 'lightgrey'])
    ax.set_title("Proportion of 'No Focused Window' ANRs")
    fig.tight_layout()
    save_chart(fig, "focused_window_anrs", img_dir)

    with open(json_path, "w", encoding="utf-8") as f:
        json.dump({
            "status": "success",
            "summary": summary}, f, indent=2, ensure_ascii=False)

    build_ppt(img_dir, ppt_path)

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True)
    parser.add_argument('--outdir', required=True)
    args = parser.parse_args()
    run_anr_analysis(args.input, args.outdir)
