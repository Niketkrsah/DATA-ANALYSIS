import sys
import os
from pptx import Presentation
from pptx.util import Inches

crash_map = {
    "Crash Frequency by Signal": "signal_freq.png",
    "Top Error Types": "error_types.png",
    "Top Backtrace Chains": "top_backtrace_chain.png",
    "Top STB Serial Numbers": "STB_Serial_NO.png",
    "Top STB Models": "STB_Model.png",
    "Faulting Libraries": "faulting_libraries.png",
    "Crashes by State": "by_state.png",
    "Memory Ranges": "memory_ranges.png",
    "Top ODU Serial Numbers": "odu_serial_top10.png",
    "Top ODU Models": "odu_model_top10.png",
    "Top Customer Product Types": "cust_product_types.png",
    "Top Plan Status Types": "plan_status.png",
    "Top GIS Building IDs": "gis_building_ids.png",
    "Top Cities by Crash": "gis_city.png",
    "Top GIS ONT Serial Numbers": "gis_ont_serials.png",
    "App version": "as5_firmware.png",
    "Top Faulting Libraries (As 5)": "as5_fault_libs.png",
    "Top Crashing Threads": "top_threads.png"
}

anr_map = {
    "Top ANR Activities": "top_anr_activities.jpg",
    "ANR by App Version": "anr_by_version.jpg",
    "ANR by State": "anr_by_state.jpg",
    "ANR by Hour": "anr_by_hour.jpg",
    "Available Memory Ranges": "mi4_available_memory_ranges.jpg",
    "Simplified ANR Subjects": "anr_subject_simplified.jpg",
    "Top 5 STB Series": "top_5_stb_series_vertical.jpg",
    "Customer Product Types": "top_5_customer_product_types.jpg",
    "No Focused Window": "focused_window_anrs.jpg"
}

if len(sys.argv) != 5:
    print("Usage: python generate_filtered_ppt.py <sessionDir> <analysisType> <filteredTitles> <pptxFilename>")
    sys.exit(1)

session_dir = sys.argv[1]
analysis_type = sys.argv[2]
filtered_titles = sys.argv[3].split(';;')
pptx_filename = sys.argv[4]

image_dir = os.path.join(session_dir, f"{analysis_type}_images")
emailppt_dir = os.path.join(session_dir, "emailppt")

# Create the directory if it doesn't exist
os.makedirs(emailppt_dir, exist_ok=True)  
pptx_path = os.path.join(emailppt_dir, pptx_filename)

#  Choose correct mapping
title_to_filename = crash_map if analysis_type == "crash" else anr_map

prs = Presentation()
blank_slide = prs.slide_layouts[6]

for title in filtered_titles:
    filename = title_to_filename.get(title)
    if filename:
        img_path = os.path.join(image_dir, filename)
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_slide)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            print(f"⚠️ Image not found: {img_path}")
    else:
        print(f"⚠️ No filename mapping found for: {title}")

prs.save(pptx_path)
