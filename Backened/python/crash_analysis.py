import os
import re
import json
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import matplotlib
import warnings
from pptx import Presentation
from pptx.util import Inches
import textwrap
import ast

# === Configs ===
matplotlib.rcParams["font.family"] = "Segoe UI Emoji"
warnings.filterwarnings("ignore", message=".*Tight layout not applied.*")
sns.set(style="whitegrid")

def run_crash_analysis(input_csv, output_dir):
    def ensure_dirs(base):
        img_dir = os.path.join(base, "crash_images")
        os.makedirs(img_dir, exist_ok=True)
        for f in os.listdir(img_dir):
            if f.endswith(".jpg") or f.endswith(".png"):
                os.remove(os.path.join(img_dir, f))
        return img_dir

    def cleanup_df(df):
        df.columns = df.columns.str.strip()
        if "Event Time" in df:
            df["Event Time"] = pd.to_datetime(df["Event Time"], errors="coerce")
        return df.ffill().infer_objects(copy=False)

    def extract(pattern, text, group=1):
        match = re.search(pattern, str(text))
        return match.group(group) if match else None

    def parse_structured_crash_logs(df):
        col = "App Crash Process Info"
        rows = []
        for log in df[col].dropna().astype(str):
            rows.append({
                "Timestamp": extract(r"Timestamp: (.*?) pid", log),
                "App": extract(r">>> (.*?) <<<", log),
                "ABI": extract(r"ABI: '(.*?)'", log),
                "Process": extract(r"name: (.*?)\s+>>>", log),
                "Signal Name": extract(r"signal \d+ \((.*?)\)", log),
                "Signal Code": int(extract(r"signal (\d+)", log) or -1),
                "Error Name": extract(r"code -?\d+ \((.*?)\)", log),
                "Error Code": int(extract(r"code (-?\d+)", log) or -1),
                "Build": extract(r"Build fingerprint: '(.*?)'", log),
                "Libraries (Backtrace)": " → ".join(re.findall(r"#\d+\s+pc .*?\s+(.*?)\s",log))
            })
        return pd.DataFrame(rows)

    def save_chart(fig, name, img_dir):
        path = os.path.join(img_dir, f"{name}.png")
        fig.savefig(path, dpi=300, bbox_inches="tight")
        plt.close(fig)
        return path

    def wrap_labels(series, width=30):
        return series.rename(lambda label: "\n".join(textwrap.wrap(label, width)))
    
    def top5_bar(colname, title, xlabel, filename, df, img_dir, summary):
        if colname in df.columns:
            vc = df[colname].value_counts().head(10)
            bar_chart(vc, title, xlabel, filename, img_dir)
            summary[title] = vc.to_dict()


    def bar_chart(series, title, xlabel, name, img_dir, horizontal=False):
        series = wrap_labels(series,width=100)

        num_items = len(series)
        max_label_len = max([len(str(label)) for label in series.index])

    # Dynamic sizing
        width = 12 if not horizontal else min(20, 6 + max_label_len * 0.2)
        height = min(10, 4 + num_items * 0.4)

        fig, ax = plt.subplots(figsize=(width, height))

        if horizontal:
            if name == "top_backtrace_chain":
                series = series[::-1]
            series.plot(kind="barh", ax=ax, color="skyblue", edgecolor="black")
            ax.invert_yaxis()
        else:
            series.plot(kind="bar", ax=ax, color="steelblue", edgecolor="black")
            plt.xticks(rotation=60, ha='right')


        ax.set_title(title)
        ax.set_xlabel(xlabel)
        ax.set_ylabel("Count")
        for p in ax.patches:
            if horizontal:
                val = p.get_width() if horizontal else p.get_height()
                ax.annotate(str(int(val)), (val + 10, p.get_y() + p.get_height() / 2),
                        va='center', fontsize=9)
            else:
                val = p.get_height()
                ax.annotate(str(int(val)), (p.get_x() + p.get_width() / 2, val + 1),
                        ha='center', fontsize=9)
        fig.tight_layout(rect=[0, 0, 1, 1])
        return save_chart(fig, name, img_dir)

    def build_ppt(img_dir, outdir):
        ppt_path = os.path.join(outdir, "crash_report.pptx")
        if os.path.exists(ppt_path): os.remove(ppt_path)

        prs = Presentation()
        blank = prs.slide_layouts[6]

        for img in sorted(os.listdir(img_dir)):
            if img.endswith(".png"):
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(os.path.join(img_dir, img), 0, 0,
                                         width=prs.slide_width, height=prs.slide_height)
        prs.save(ppt_path)

    # === Load and Prepare Data ===
    df = pd.read_csv(input_csv)
    df = cleanup_df(df)
    img_dir = ensure_dirs(output_dir)
    summary_path = os.path.join(output_dir, "crash_summary.json")
    if os.path.exists(summary_path): os.remove(summary_path)
    summary = {}


    # === Parse Crash Process Logs ===
    parsed = parse_structured_crash_logs(df)

    if not parsed.empty:
        if "Signal Name" in parsed:
            vc = parsed["Signal Name"].value_counts()
            bar_chart(vc, "Crash Frequency by Signal", "Signal", "signal_freq", img_dir)
            summary["Crash Frequency by Signal"] = vc.to_dict()

        if "Error Name" in parsed:
            vc = parsed["Error Name"].value_counts()
            bar_chart(vc, "Top Error Types", "Error", "error_types", img_dir)
            summary["Top Error Types"] = vc.to_dict()

        if "Libraries (Backtrace)" in parsed:
            vc = parsed["Libraries (Backtrace)"].value_counts().head(5)
            bar_chart(vc, "Top Backtrace Chains", "Backtrace", "top_backtrace_chain", img_dir, horizontal=True)
            summary["Top Backtrace Chains"] = vc.to_dict()

        parsed["Timestamp"] = pd.to_datetime(parsed["Timestamp"], errors="coerce", utc=True)
        # parsed.to_csv(os.path.join(output_dir, "parsed_crash_info.csv"), index=False)

    # === STB Serial and Model Information ===
    if "Di 8" in df.columns:
        vc = df["Di 8"].value_counts().nlargest(10)
        bar_chart(vc, "Top 10 STB Serial Numbers", "STB Serial No", "STB_Serial_NO", img_dir)
        summary["Top STB Serial Numbers"] = vc.to_dict()

    if "STB model" in df.columns:
        vc = df["STB model"].value_counts().nlargest(10)
        bar_chart(vc, "Top 10 STB Models", "STB model", "STB_Model", img_dir)
        summary["Top STB Models"] = vc.to_dict()

    # === Libraries involved ===
    def parse_libs(log):
        match = re.search(r'/lib/.*/(lib\w+\.so)', str(log))
        return match.group(1) if match else None

    libs = df["App Crash Process Info"].dropna().apply(parse_libs).dropna()
    if not libs.empty:
        vc = libs.value_counts().nlargest(10)
        fig, ax = plt.subplots(figsize=(10, 5))
        sns.countplot(y=libs, order=vc.index, ax=ax, palette="crest")


        for container in ax.containers:
            for bar in container:
                width = bar.get_width()
                ax.annotate(f'{int(width)}', xy=(width + 5, bar.get_y() + bar.get_height() / 2),
                    va='center', fontsize=9)

        ax.set_title("Top Faulting Libraries")
        ax.set_xlabel("Crash Count")
        ax.set_ylabel("Library")
        fig.tight_layout()

        save_chart(fig, "faulting_libraries", img_dir)
        summary["Faulting Libraries"] = vc.to_dict()

    # === Crashes by State ===
    if "State" in df.columns:
        vc = df["State"].value_counts()
        bar_chart(vc, "Crashes by State", "State", "by_state", img_dir)
        summary["Crashes by State"] = vc.to_dict()
    # Extract Available Memory from Mi 4 column
    if "Mi 4" in df.columns:
        def extract_available_memory(val):
            try:
                return json.loads(val).get("available", 0)
            except:
                return 0

        df["Available Memory"] = df["Mi 4"].apply(extract_available_memory)
        bins = [100, 400, 800, 1000, 1200, 1400, 2000]
        labels = ["100-400", "400-800", "800-1000", "1000-1200", "1200-1400", "1600-2000"]
        df["Memory Range"] = pd.cut(df["Available Memory"], bins=bins, labels=labels, include_lowest=True)

        vc = df["Memory Range"].value_counts().sort_index()
        bar_chart(vc, "Memory Availability Distribution", "Range (MB)", "memory_ranges", img_dir)
        summary["Memory Ranges"] = vc.to_dict()
    # === Acs Odu Serial Number & Model ===
    if "Acs Odu Serial Number" in df.columns:
        vc = df["Acs Odu Serial Number"].value_counts().nlargest(10)
        bar_chart(vc, "Top ODU Serial Numbers", "ODU Serial", "odu_serial_top10", img_dir)
        summary["Top ODU Serial Numbers"] = vc.to_dict()

    if "Acs Odu Model" in df.columns:
        vc = df["Acs Odu Model"].value_counts().nlargest(10)
        bar_chart(vc, "Top ODU Models", "ODU Model", "odu_model_top10", img_dir)
        summary["Top ODU Models"] = vc.to_dict()

    # === Customer, Plan, GIS Metadata ===
    top5_bar("D Customer Product Type", "Top Customer Product Types", "Product Type", "cust_product_types", df, img_dir, summary)
    top5_bar("D Plan Status", "Top Plan Status Types", "Plan Status", "plan_status", df, img_dir, summary)
    top5_bar("Gis Building Id", "Top GIS Building IDs", "Building ID", "gis_building_ids", df, img_dir, summary)
    top5_bar("Gis City", "Top Cities by Crash", "City", "gis_city", df, img_dir, summary)
    top5_bar("Gis Ont Serial Number", "Top GIS ONT Serial Numbers", "ONT Serial", "gis_ont_serials", df, img_dir, summary)

    def parse_as5_column(df):
        parsed_data = []
    
        if "As 5" not in df.columns:
            return pd.DataFrame()

        for raw in df["As 5"].dropna():
            try:
                entry_list = ast.literal_eval(raw)
                if isinstance(entry_list, list) and entry_list:
                    item = entry_list[0]
                    ps21 = item.get("PS21")
                    ps22 = item.get("PS22", "")
                    xps46 = item.get("XPS46")
                    xps911 = item.get("XPS911")

                    # Timestamp
                    timestamp = re.search(r'Timestamp:\s*(.*?)\s+pid', ps22)
                    timestamp = timestamp.group(1) if timestamp else None

                    # First Faulting Library
                    fault_lib = None
                    backtrace = re.findall(r'#\d+\s+pc .*?/(lib\w+\.so)', ps22)
                    if backtrace:
                        fault_lib = backtrace[0]

                    thread_name = None
                    thread_match = re.search(r'name:\s*(.*?)\s+>>>', ps22)
                    if thread_match:
                         thread_name = thread_match.group(1).strip()

                    parsed_data.append({
                        "App": ps21,
                        "Timestamp": timestamp,
                        "Firmware": xps46,
                        "Subtype": xps911,
                        "Top Lib": fault_lib,
                        "Thread Name": thread_name  
                    })

            except Exception as e:
                continue

        return pd.DataFrame(parsed_data)
    parsed_as5 = parse_as5_column(df)

    if not parsed_as5.empty:
            # ───── Firmware Distribution ─────
            vc_firmware = parsed_as5["Firmware"].value_counts().head(10)
            bar_chart(vc_firmware, "App Version", "Firmware", "as5_firmware", img_dir)
            summary["App version"] = vc_firmware.to_dict()

            # ───── Top Faulting Libraries ─────
            vc_libs = parsed_as5["Top Lib"].dropna().value_counts().head(10)
            bar_chart(vc_libs, "Top Faulting Libraries (As 5)", "Library", "as5_fault_libs", img_dir)
            summary["Top Faulting Libraries (As 5)"] = vc_libs.to_dict()

            
            vc_name = parsed_as5["Thread Name"].value_counts().head(10)
            bar_chart(vc_name, "Top Crashing Threads", "Thread Name", "top_threads", img_dir)
            summary["Top Crashing Threads"]= vc_name.to_dict()


# ─────────────── Export Outputs ───────────────

    # === Save JSON + PPT ===
    with open(summary_path, "w", encoding="utf-8") as f:
        json.dump({
            "status": "success",
        "summary": summary },f, indent=2, ensure_ascii=False)

    # Build PowerPoint with charts
    build_ppt(img_dir, output_dir)


# === CLI Support ===
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True)
    parser.add_argument('--outdir', required=True)
    args = parser.parse_args()
    run_crash_analysis(args.input, args.outdir)
